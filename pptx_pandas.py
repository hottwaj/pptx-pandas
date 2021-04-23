from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.exc import PackageNotFoundError
import pptx
import numbers
import pandas
from itertools import zip_longest
from shutil import copyfile
from typing import Union, Protocol

class PrettyPandas(Protocol):
    def get_formatted_df(self) -> pandas.DataFrame: pass

Table = Union[pandas.DataFrame, PrettyPandas]

class PlotlyChartBundle(Protocol):
    def write_image(self, filename: str, scale: float, *args, **kwargs): pass

class MatplotlibFigure(Protocol):
    def savefig(filename: str, *args, **kwargs): pass

class MatplotlibPlot(Protocol):
    figure: MatplotlibFigure

Chart = Union[PlotlyChartBundle, MatplotlibPlot]

def grouped(iterable, n):
    "group a sequence of objects into a sequence of tuples each containing 'n' objects"
    "s -> (s0,s1,s2,...sn-1), (sn,sn+1,sn+2,...s2n-1), (s2n,s2n+1,s2n+2,...s3n-1), ..."
    #from http://stackoverflow.com/a/5389547/1280629
    return zip_longest(*[iter(iterable)]*n)

class NoExistingSlideFoundError(RuntimeError): pass

DIST_METRIC = Inches
DEFAULT_POSITIONS = {
    'content': [5.2, 4.1, 5, 2.7],
    'index_width': 0.7, 
    'col_width': 0.6,
    'table_row_height': 0.2,
    'single_item_top': 1.7, 
    'single_item_left': 2.2,
    'multi_item_top': 1.4, 
    'multi_item_left': 1.5,
    'horizontal_gap': 0.7,
    'vertical_gap': 0.5
}
    
class PresentationWriter():
    def __init__(self, pptx_file: str, pptx_title: str, pptx_template: str = None, 
                 default_overwrite_if_present: bool = False,
                 default_overwrite_only: bool = False,
                 default_position_overrides: dict = {},
                 default_table_font_attrs: dict = {'size': Pt(7), 'name': 'Calibri'},
                 default_caption_font_attrs: Optional[dict] = None,
                 default_include_internal_cell_borders: bool = True,
                 default_border_kwargs: dict = {},
                 default_slide_layout: str = 'Normal Page',
                 default_remove_empty_text_boxes: bool = False,
                 default_auto_position_elements: bool = False,
                 save_every: Optional[int] = 1):
        self.pptx_file = pptx_file
        self.pptx_title = pptx_title
        
        if pptx_template is not None:
            copyfile(pptx_template, pptx_file)
        
        try:
            self.presentation = Presentation(self.pptx_file) #open existing file if present
        except PackageNotFoundError as ex:
            if pptx_template is None:
                self.presentation = Presentation() #open a blank presentation
            else:
                raise PackageNotFoundError(f'Could not open file file "{pptx_file}" after copying template file "{pptx_template}" onto it') from ex # re-raise the exception
        core_props = self.presentation.core_properties
        core_props.title = pptx_title
        core_props.keywords = ''
        
        self.default_overwrite_if_present = default_overwrite_if_present
        self.default_overwrite_only = default_overwrite_only
        self.default_positions = {**DEFAULT_POSITIONS, **default_position_overrides}
        self.default_table_font_attrs = default_table_font_attrs
        self.default_caption_font_attrs = (default_caption_font_attrs
                                           if default_caption_font_attrs is not None
                                           else {**{'bold': True}, **default_table_font_attrs})
        self.default_include_internal_cell_borders = default_include_internal_cell_borders
        self.default_border_kwargs = default_border_kwargs
        self.default_slide_layout = default_slide_layout
        self.default_remove_empty_text_boxes = default_remove_empty_text_boxes
        self.default_auto_position_elements = default_auto_position_elements
        self.save_every = save_every
        self.slide_count = 0
        
    def save_presentation(self):
        self.presentation.save(self.pptx_file)

    def _new_slide_from_template(self, slide_layout_name: Optional[str] = None):
        prs = self.presentation
        slide_layout_name = slide_layout_name or self.default_slide_layout
        slide_layout = prs.slide_layouts.get_by_name(slide_layout_name)
        if slide_layout is None:
            possible_layouts = ', '.join([sl.name for sl in prs.slide_layouts])
            raise ValueError(f'slide_layout="{slide_layout_name}" is not a valid layout in this presentation file. Options are: {possible_layouts}')
        return prs.slides.add_slide(slide_layout)

    @classmethod
    def _add_chart_to_slide(cls, chart: Chart, slide: pptx.slide.Slide, left: float, top: float, tmp_filename: str):
        if hasattr(chart, 'width'):
            chart_width = chart.width / 800.0 * 5.0
            chart_height = chart.height / 800.0 * 5.0
        else:
            chart_width = chart.figure.get_figwidth()  # already in inches
            chart_height = chart.figure.get_figheight()

        chart_to_file(chart, tmp_filename)

        return slide.shapes.add_picture(tmp_filename,
                                        left = left,
                                        top = top,
                                        width = Inches(chart_width),
                                        height = Inches(chart_height))

    def _add_table_to_slide(self, table: Table, slide: pptx.slide.Slide, left: float, top: float,
                            positions: dict,
                            table_font_attrs: Optional[dict] = None,
                            border_kwargs: Optional[dict] = None,
                            include_internal_cell_borders: Optional[bool] = None,
                            caption_font_attrs: Optional[dict] = None):
        table_df = table_to_dataframe(table)
        col_widths = [positions['index_width']] + [positions['col_width']] * len(table_df.columns)
        table_shape = create_pptx_table(slide, table_df,
                                        left = left,
                                        top = top,
                                        col_width = col_widths,
                                        row_height = positions['table_row_height'],
                                        font_attrs = table_font_attrs or self.default_table_font_attrs,
                                        border_kwargs = border_kwargs or self.default_border_kwargs,
                                        include_internal_cell_borders = include_internal_cell_borders
                                                                        if include_internal_cell_borders is not None
                                                                        else self.default_include_internal_cell_borders)
        table_caption = getattr(table, 'caption')
        if table_caption:
            caption_font_attrs = caption_font_attrs or self.default_caption_font_attrs
            caption_box = slide.shapes.add_textbox(left = left,
                                                   top = top - caption_font_attrs['size'] * 2,
                                                   width = table_shape.width,
                                                   height = caption_font_attrs['size'])
            set_cell_text(caption_box, table_caption)
            set_cell_font_attrs(caption_box, **caption_font_attrs)
            table_shape._pptx_pandas_caption = caption_box

        return table_shape

    positions_deprecations = {'single_table_margin_top': 'single_item_top',
                              'single_table_margin_left': 'single_item_left',
                              'single_item_margin_top': 'single_item_top',
                              'single_item_margin_left': 'single_item_left',
                              'multi_item_margin_top': 'multi_item_top',
                              'multi_item_margin_left': 'multi_item_left',
                              'charts_horizontal_gap': 'horizontal_gap',
                              'charts_vertical_gap': 'vertical_gap',}

    @classmethod
    def check_positions_settings(cls, positions):
        bad_position_setting_list = []
        for k in positions:
            if k in cls.positions_deprecations:
                bad_position_setting_list.append(f'{k} -> {cls.positions_deprecations[k]}')
            elif k not in DEFAULT_POSITIONS:
                bad_position_setting_list.append(f'{k} is not a valid position setting')
        if bad_position_setting_list:
            raise ValueError('Deprecated or invalid positions settings provided:\n'
                             + '\n'.join(bad_position_setting_list))

    write_slide_deprecated_kwargs = {'charts': 'elements',
                                     'tables': 'elements',
                                     'charts_and_tables': 'elements',
                                     'charts_per_row': 'elements_per_row',
                                     'tables_per_row': 'elements_per_row',
                                     'auto_position_charts_and_tables': 'auto_position_elements',}

    def write_slide(self, title: str,
                    elements: Union[list[list[Union[Chart, Table]]], list[Union[Chart, Table]]],
                    overwrite_if_present: Optional[bool] = None,
                    overwrite_only: Optional[bool] = None,
                    elements_per_row: Optional[int] = None,
                    position_overrides: dict = {},
                    auto_position_elements: Optional[bool] = False,
                    remove_empty_text_boxes: Optional[bool] = None,
                    table_font_attrs: Optional[dict] = None,
                    caption_font_attrs: Optional[dict] = None,
                    include_internal_cell_borders: Optional[bool] = None,
                    border_kwargs: Optional[dict] = None,
                    slide_layout_name: Optional[str] = None,
                    **kwargs) -> pptx.slide.Slide:

        deprecated_kwargs_used = []
        not_known_kwargs = []
        for k in kwargs:
            new_k = self.write_slide_deprecated_kwargs.get(k)
            if new_k is not None:
                deprecated_kwargs_used.append(f'{k} -> {new_k}')
            else:
                not_known_kwargs.append(k)
        if deprecated_kwargs_used:
            raise ValueError(
                'PresentationWriter.write_slide: following args are deprecated and should be changed as follows:\n' + '\n'.join(
                    deprecated_kwargs_used))
        if not_known_kwargs:
            raise ValueError(
                'PresentationWriter.write_slide: following kwargs are not known (possibly deprecated): ' + ', '.join(
                    not_known_kwargs))

        if elements:
            if not isinstance(elements[0], list):
                if elements_per_row is None:
                    # wrap the elements in a single row
                    elements = [elements]
                else:
                    # use elements_per_row to wrap the elements list into grouped rows
                    elements = grouped(elements, elements_per_row)

        positions = {**self.default_positions, **position_overrides}
        self.check_positions_settings(positions)

        overwrite_if_present = overwrite_if_present if overwrite_if_present is not None else self.default_overwrite_if_present
        overwrite_only = overwrite_only if overwrite_only is not None else self.default_overwrite_only
        if overwrite_if_present or overwrite_only:
            try:
                self.overwrite_pptx(title, elements)
            except NoExistingSlideFoundError:
                if overwrite_only:
                    raise
            else:
                return

        slide = self._new_slide_from_template(slide_layout_name)

        content, subtitle, title_shape, *others = [s for s in slide.shapes if s.has_text_frame]
        title_shape.text = self.pptx_title
        subtitle.text = title

        content.left, content.top, content.width, content.height = (DIST_METRIC(x) for x in positions['content'])

        shapes_added = []
        table_shapes = []
        num_charts = 0
        num_elements = sum([len(row) for row in elements])
        initial_left = DIST_METRIC(positions['multi_item_left']
                                   if num_elements > 1
                                   else positions['single_item_left'])
        curr_top = DIST_METRIC(positions['multi_item_top']
                               if num_elements > 1
                               else positions['single_item_top'])
        for row in elements:
            curr_left = initial_left
            shape_row = []
            row_height = 0
            for elem in row:
                if is_table_instance(elem):
                    shape = self._add_table_to_slide(elem, slide,
                                                     left = curr_left,
                                                     top = curr_top,
                                                     positions = positions,
                                                     table_font_attrs = table_font_attrs,
                                                     border_kwargs = border_kwargs,
                                                     include_internal_cell_borders = include_internal_cell_borders,
                                                     caption_font_attrs = caption_font_attrs)
                    table_shapes.append(shape)
                elif is_chart_instance(elem):
                    num_charts += 1
                    shape = self._add_chart_to_slide(elem, slide,
                                                     left = curr_left,
                                                     top = curr_top,
                                                     tmp_filename = f'pptx-img-{num_charts}.png')
                else:
                    raise ValueError(f'Cannot render elements of type {type(elem)} to pptx')
                curr_left += shape.width + DIST_METRIC(positions['horizontal_gap'])
                row_height = max(row_height, shape.height)
                shape_row.append(shape)

            curr_top += row_height + DIST_METRIC(positions['vertical_gap'])
            if shape_row:
                shapes_added.append(shape_row)

        auto_position_elements = auto_position_elements if auto_position_elements is not None else self.default_auto_position_elements
        if auto_position_elements:
            content_area_top = subtitle.top + subtitle.height
            content_area_height = self.presentation.slide_height - content_area_top
            if len(others) == 1:
                footnotes = others[0] # assume remaining text shape is cell footnotes
                content_area_height -= self.presentation.slide_height - footnotes.top
            # calculate shift from slide centre needed to put shapes in centre between bottom of subtitle and top of footnotes (if present)
            centre_y_shift = content_area_height/2 + content_area_top - self.presentation.slide_height/2
            self.auto_position_shapes(shapes_added,
                                      DIST_METRIC(positions['horizontal_gap']),
                                      DIST_METRIC(positions['vertical_gap']),
                                      centre_y_shift)
            for table_shape in table_shapes:
                caption_box = getattr(table_shape, '_pptx_pandas_caption', None)
                if caption_box:
                    caption_box.left = table_shape.left
                    caption_box.top = table_shape.top - caption_box.height*2

        remove_empty_text_boxes = remove_empty_text_boxes if remove_empty_text_boxes is not None else self.default_remove_empty_text_boxes
        if remove_empty_text_boxes:
            self.remove_empty_text_boxes(slide)
        self.slide_count += 1
        if self.save_every is not None and (self.slide_count % self.save_every) == 0:
            self.save_presentation()
        return slide

    def auto_position_shapes(self, shapes_by_row: list[list[pptx.shapes.base.BaseShape]],
                             horizontal_gap: int, vertical_gap: int, centre_y_shift: int = 0):
        shapes_by_col = list(zip_longest(*shapes_by_row))
        col_widths = [max([s.width for s in col if s is not None]) for col in shapes_by_col]
        row_heights = [max([s.height for s in row]) for row in shapes_by_row]
        total_width = sum(col_widths)
        total_height = sum(row_heights)

        total_height_inc_margins = total_height + (len(row_heights)-1)*vertical_gap
        top_pos = (self.presentation.slide_height - total_height_inc_margins) / 2
        for cell_height, row in zip(row_heights, shapes_by_row):
            if len(row) < len(col_widths):
                # row that has fewer items than other rows
                # with current layout logic this can happen if we have a final ragged row e.g. rows of 3, 3, 3, 1)
                # can also happen if different charts_per_row and tables_per_row values are used.  might get odd results in that situation
                cell_widths = [s.width for s in row]
                row_width = sum(cell_widths)
            else:
                cell_widths = col_widths
                row_width = total_width
            row_width += (len(row)-1) * horizontal_gap

            # calculate left position to centre the cell on the page, given row width
            left_pos = (self.presentation.slide_width - row_width) / 2
            for shape, width in zip(row, cell_widths):
                # now position the shape in the centre of containing cell
                shape.top = int(top_pos + (cell_height - shape.height) / 2 + centre_y_shift)
                shape.left = int(left_pos + (width - shape.width) / 2)
                left_pos += width + horizontal_gap

            top_pos += cell_height

    @classmethod
    def remove_empty_text_boxes(cls, slide):
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.text:
                elem = shape.element
                elem.getparent().remove(elem)
    
    def overwrite_pptx(self, slide_title,
                       elements: list[list[Union[Chart, Table]]]):
        prs = self.presentation

        for i, slide in enumerate(prs.slides):
            text_boxes = [s for s in slide.shapes if s.has_text_frame]

            matched_title = False
            for tb in text_boxes:
                if tb.text == slide_title:
                    matched_title = True
                    break

            if matched_title:
                break

        if matched_title:
            found_shapes = []
            for s in slide.shapes:
                if isinstance(s, pptx.shapes.picture.Picture):
                    found_shapes.append(s)
                elif isinstance(s, pptx.shapes.graphfrm.GraphicFrame) and s.has_table:
                    found_shapes.append(s)

            flattened_elements = []
            for row in elements:
                flattened_elements.extend(row)

            for i, (elem, shape) in enumerate(zip(flattened_elements, found_shapes)):
                if elem == chart and shape == chart:
                    img_file = 'pptx-img-%d.png' % i
                    chart_to_file(chart, img_file)

                    # Replace image:
                    picture = shape
                    with open(img_file, 'rb') as f:
                        imgBlob = f.read()
                    imgRID = picture._pic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                    imgPart = slide.part.related_parts[imgRID]
                    imgPart._blob = imgBlob
                elif elem == table and shape == table:
                    table_df = table_to_dataframe(table)
                    write_pptx_dataframe(table_df, shape.table, overwrite_formatting = False)

            self.save_presentation()
        else:
            raise NoExistingSlideFoundError('No existing slide titled "%s"' % slide_title)

def is_table_instance(obj: Any) -> bool:
    return isinstance(obj, pandas.DataFrame) or hasattr(obj, 'get_formatted_df')

def is_chart_instance(obj: Any) -> bool:
    return hasattr(obj, 'write_image') or hasattr(getattr(obj, 'figure', None), 'savefig')

def table_to_dataframe(table: Table) -> pandas.DataFrame:
    if isinstance(table, pandas.DataFrame):
        table_df = table.copy()
    elif hasattr(table, 'get_formatted_df'):
        table_df = table.get_formatted_df()
    else:
        return ValueError('table_to_dataframe(): Given table object is not a DataFrame or PrettyPandas object')
    table_df.columns = [s.replace('<br>', '\n')
                        if isinstance(s, str) else s
                        for s in table_df.columns]
    return table_df

empty_text_set = {'', None}
def set_cell_text(cell, text, overwrite_formatting = True):
    if text in empty_text_set:
        text = "\u00A0" # unicode nbsp - needed to fill empty cells as otherwise formatting is not applied by PPT
    if overwrite_formatting:
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = text
    else:
        p = cell.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].text = text
        else:
            r = p.add_run()
            r.text = text
        
def set_cell_font_attrs(cell, **kwargs):
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            for k, v in kwargs.items():
                if k == 'color_rgb':
                    r.font.color.rgb = v
                else:
                    setattr(r.font, k, v)
                
def format_cell_text(val, float_format = '{:.0f}', int_format = '{:d}'):
    if isinstance(val, numbers.Integral):
        return int_format.format(val)
    elif isinstance(val, numbers.Real):
        return float_format.format(val)
    else:
        return str(val)
    
def set_cell_borders(cell, border_color="4f81bd", border_width='6350', border_scheme_color = 'accent1', borders = 'LRTB'):
    """ Hack function to enable the setting of border width and border color"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for border in borders:
        lnL = SubElement(tcPr, 'a:ln' + border, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        lnL_solidFill = SubElement(lnL, 'a:solidFill')
        if border_scheme_color is not None:
            lnL_srgbClr = SubElement(lnL_solidFill, 'a:schemeClr', val=border_scheme_color)
        else:
            lnL_srgbClr = SubElement(lnL_solidFill, 'a:srgbClr', val=border_color)
        lnL_prstDash = SubElement(lnL, 'a:prstDash', val='solid')
        lnL_round_ = SubElement(lnL, 'a:round')
        lnL_headEnd = SubElement(lnL, 'a:headEnd', type='none', w='med', len='med')
        lnL_tailEnd = SubElement(lnL, 'a:tailEnd', type='none', w='med', len='med')
    cell.fill.background()
    
def get_index_numlevels(pd_index):
    if isinstance(pd_index, pandas.MultiIndex):
        return len(pd_index.levels)
    else:
        return 1
        
def write_pptx_dataframe(dataframe, pptx_table, col_width = 1.0, format_opts = {}, font_attrs = {'size': Pt(8), 'name': 'Calibri'},
                         header_font_attrs = {'bold': True, 'color_rgb': RGBColor(34, 64, 97)},
                         border_kwargs = {},
                         include_internal_cell_borders = True,
                         overwrite_formatting = True):
    num_rows, num_cols = dataframe.shape
    if isinstance(dataframe.index, pandas.MultiIndex):
        raise RuntimeError('Cannot yet cope with MultiIndex in rows')
    num_indexes = 1
    
    num_header_rows = get_index_numlevels(dataframe.columns)
        
    if (num_rows + num_header_rows) != len(pptx_table.rows):
        raise RuntimeError('Need %d rows but PPTX table has %d'
                           % (num_rows + num_header_rows, len(pptx_table.rows)))
    if (num_cols + num_indexes) != len(pptx_table.columns):
        raise RuntimeError('Need %d columns but PPTX table has %d'
                           % (num_cols + num_indexes, len(pptx_table.columns)))
        
    header_font_attrs = dict(header_font_attrs, **font_attrs)
    last_col = len(dataframe.columns.values)-1
    for i in range(num_header_rows):
        #headers
        prev_header = no_prev_header = '##special missing value'
        first_merged_cell = 0
        mergeable_cell_count = 0
        for c, header_name in enumerate(dataframe.columns.values):
            col_name = header_name[i] if num_header_rows > 1 else header_name
            if prev_header == no_prev_header or prev_header != col_name:
                if mergeable_cell_count > 0:
                    pptx_table.cell(i, first_merged_cell + num_indexes).merge(pptx_table.cell(i, first_merged_cell + mergeable_cell_count + num_indexes))
                prev_header = col_name
                first_merged_cell = c
                mergeable_cell_count = 0
                
                cell = pptx_table.cell(i, c + num_indexes)
                set_cell_text(cell, 
                              format_cell_text(col_name, **format_opts), 
                              overwrite_formatting = overwrite_formatting)
                if overwrite_formatting:
                    set_cell_font_attrs(cell, **header_font_attrs)
                    if include_internal_cell_borders:
                        borders = 'LRTB'
                    else:
                        borders = '' + ('L' if c == 0 else '') + ('R' if c == last_col else '') + ('T' if i == 0 else '')
                    set_cell_borders(cell, borders = borders, **border_kwargs)
            else:
                mergeable_cell_count += 1
        if mergeable_cell_count > 0:
            pptx_table.cell(i, first_merged_cell + num_indexes).merge(pptx_table.cell(i, first_merged_cell + mergeable_cell_count + num_indexes))
    
    last_row = num_rows-1
    last_col = num_cols-1
    for c in range(num_cols):
        #set column widths
        if isinstance(col_width, numbers.Number):
            w = DIST_METRIC(col_width)
        else:
            w = DIST_METRIC(col_width[c + 1])
        if overwrite_formatting:
            pptx_table.columns[c + num_indexes].width = w

        #body cells
        for r in range(num_rows):
            cell = pptx_table.cell(r + num_header_rows, c + num_indexes)
            set_cell_text(cell, 
                          format_cell_text(dataframe.iloc[r, c], **format_opts), 
                          overwrite_formatting = overwrite_formatting)
            if overwrite_formatting:
                set_cell_font_attrs(cell, **font_attrs)
                if include_internal_cell_borders:
                    borders = 'LRTB'
                else:
                    borders = '' + ('R' if c == last_col else '') + ('B' if r == last_row else '')
                set_cell_borders(cell, borders = borders, **border_kwargs)

    #index
    for r in range(num_rows):
        cell = pptx_table.cell(r + num_header_rows, 0)
        set_cell_text(cell, 
                      format_cell_text(dataframe.index[r], **format_opts), 
                      overwrite_formatting = overwrite_formatting)

        if overwrite_formatting:
            set_cell_font_attrs(cell, **header_font_attrs)
            if include_internal_cell_borders:
                borders = 'LRTB'
            else:
                borders = 'L' + ('B' if r == last_row else '')
            set_cell_borders(cell, borders = borders, **border_kwargs)
    
    if overwrite_formatting:
        pptx_table.columns[0].width = DIST_METRIC(col_width if isinstance(col_width, numbers.Number) else col_width[0])
    
    #index name
    for i in range(num_header_rows):
        cell = pptx_table.cell(i, 0)
        set_cell_text(cell, 
                      format_cell_text(dataframe.index.name if dataframe.index.name is not None else '', **format_opts), 
                      overwrite_formatting = overwrite_formatting)
        if overwrite_formatting:
            set_cell_font_attrs(cell, **header_font_attrs)
            if include_internal_cell_borders:
                borders = 'LRTB'
            else:
                borders = 'L' + ('T' if i == 0 else '')
            set_cell_borders(cell, borders = borders, **border_kwargs)
        
def create_pptx_table(pptx_slide, dataframe, left, top, col_width, row_height, include_internal_cell_borders, **write_kwargs):
    num_rows, num_cols = dataframe.shape
    if isinstance(dataframe.index, pandas.MultiIndex):
        raise RuntimeError('Cannot yet cope with MultiIndex rows')
    num_indexes = 1

    num_header_rows = get_index_numlevels(dataframe.columns)

    width = DIST_METRIC(col_width * num_cols if isinstance(col_width, numbers.Number) else sum(col_width))
    height = DIST_METRIC(row_height * num_rows)

    table_shape = pptx_slide.shapes.add_table(num_rows + num_header_rows, num_cols + num_indexes, 
                             left, top,
                             width, height)
    
    table = table_shape.table
    write_pptx_dataframe(dataframe, table, col_width = col_width, include_internal_cell_borders = include_internal_cell_borders, **write_kwargs)
    return table_shape
    
def chart_to_file(chart_obj, img_file: str):
    if hasattr(chart_obj, 'write_image'):
        #ply_pd PlotlyChartBundle plots
        chart_obj.write_image(img_file, scale=2, width = chart_obj.width, height = chart_obj.height)
    else:
        #matplotlib plots...
        chart_obj.figure.savefig(img_file, dpi = 300, bbox_inches = 'tight')

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element
